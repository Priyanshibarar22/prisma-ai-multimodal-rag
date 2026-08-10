import fitz  # PyMuPDF, for PDFs
import docx # python-docx, for Word docs
from pptx import Presentation
import os

def extract_pdf(path):
    try:
        doc = fitz.open(path)
        text = ""
        for page_num, page in enumerate(doc):
            text += f"\n[Page {page_num+1}]\n" + page.get_text()
        doc.close()
        return text, None
    except Exception as e:
        return None, f"Could not read PDF: {e}"

def extract_docx(path):
    try:
        doc = docx.Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return text, None
    except Exception as e:
        return None, f"Could not read DOCX: {e}"

def extract_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return text, None
    except Exception as e:
        return None, f"Could not read TXT: {e}"
    
def extract_pptx(path):
    try:
        prs = Presentation(path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n[Slide {i+1}]\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n"
        return text, None
    except Exception as e:
        return None, f"Could not read PPTX: {e}"    

def extract_any(path):
    ext = path.split(".")[-1].lower()

    if ext == "pdf":
        text, error = extract_pdf(path)
    elif ext == "docx":
        text, error = extract_docx(path)
    elif ext == "txt":
        text, error = extract_txt(path)
    elif ext == "pptx":
        text, error = extract_pptx(path)
    else:
        return None, f"Unsupported file type: .{ext}"

    if error:
        return None, error

    if not text or not text.strip():
        return None, "File appears to be empty or has no extractable text."

    return text, None
